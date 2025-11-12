import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
from pathlib import Path
import json
import asyncio
from agents import Agent, Runner, WebSearchTool, function_tool, ModelSettings
from typing import Dict, Any, List, Optional
from pydantic import BaseModel, Field
import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv

load_dotenv()

# Importar herramientas de mcp_tools con manejo de errores
try:
    from mcp_tools import tavily_search, wikipedia_search, duckduckgo_search
except ImportError as e:
    import warnings
    warnings.warn(f"No se pudieron importar algunas herramientas de mcp_tools: {e}")
    # Funciones de respaldo
    def tavily_search(query: str, search_depth: str = "basic") -> str:
        return f"Error: Tavily no disponible. Query: {query}"
    def wikipedia_search(query: str) -> Dict[str, Any]:
        return {"success": False, "error": "Wikipedia no disponible"}
    def duckduckgo_search(query: str) -> Dict[str, Any]:
        return {"success": False, "error": "DuckDuckGo no disponible"}

st.set_page_config(page_title="Análisis de Oportunidad Académica", layout="wide")

# -------------------------------------------------------------
# Modelos de datos para los agentes
# -------------------------------------------------------------

class ContextoPrograma(BaseModel):
    """Modelo para el contexto del programa organizado en JSON"""
    nombre_programa: str = Field(..., description="Nombre del programa a analizar")
    descripcion: str = Field(..., description="Descripción del programa")
    nivel_academico: str = Field(..., description="Nivel académico (pregrado, maestría, doctorado)")
    palabras_clave: List[str] = Field(default_factory=list, description="Palabras clave para búsqueda")
    contexto_geografico: str = Field(default="Colombia", description="Contexto geográfico de interés")
    objetivos_busqueda: List[str] = Field(default_factory=list, description="Objetivos de la búsqueda")

class ProgramaEncontrado(BaseModel):
    """Modelo para programas encontrados en la búsqueda"""
    nombre: str = Field(..., description="Nombre del programa")
    universidad: str = Field(..., description="Nombre de la universidad")
    pais: str = Field(..., description="País donde se ofrece")
    url: Optional[str] = Field(None, description="URL del programa")
    descripcion: Optional[str] = Field(None, description="Descripción del programa")
    nivel: Optional[str] = Field(None, description="Nivel académico")
    cursos_representativos: List[str] = Field(default_factory=list, description="Cursos representativos")
    fuente: str = Field(..., description="Fuente de la información")

class ResumenPrograma(BaseModel):
    """Modelo para el resumen y score de relación"""
    programa_original: str
    programas_similares: List[ProgramaEncontrado]
    score_relacion: Dict[str, float] = Field(default_factory=dict, description="Score de relación por programa")
    resumen_comparativo: str = Field(..., description="Resumen comparativo de los programas")
    tendencias_nombres: List[str] = Field(default_factory=list, description="Tendencias en nombres encontrados")

class ReporteFinal(BaseModel):
    """Modelo para el reporte final agregado"""
    contexto: ContextoPrograma
    resumen: ResumenPrograma
    variables_cuantitativas: Dict[str, Any] = Field(default_factory=dict)
    recomendaciones: List[str] = Field(default_factory=list)

# -------------------------------------------------------------
# Funciones de apoyo
# -------------------------------------------------------------

def cargar_datos():
    """Carga los datasets de SNIES desde la fuente pública."""
    maestro = pd.read_parquet('https://robertohincapie.com/data/snies/MAESTRO.parquet')
    oferta = pd.read_parquet('https://robertohincapie.com/data/snies/OFERTA.parquet')
    programas = pd.read_parquet('https://robertohincapie.com/data/snies/PROGRAMAS.parquet')
    ies = pd.read_parquet('https://robertohincapie.com/data/snies/IES.parquet')
    return maestro, oferta, programas, ies


def analizar_programa(programa_nombre: str):
    """Ejecuta el flujo de análisis SNIES para el programa indicado."""
    maestro, oferta, programas, ies = cargar_datos()
    requerido = set(programa_nombre.lower().split())
    programa = set(programa_nombre.lower().split())
    n = len(programa)

    equivalentes = []
    for prg in programas['PROGRAMA_ACADEMICO'].unique():
        prg2 = str(prg).lower().split()
        indice = len(set(programa).intersection(prg2)) / len(programa)
        if indice >= (n - 1) / n and len(requerido.intersection(prg2)) == len(requerido):
            equivalentes.append(prg)

    programas2 = programas[programas['PROGRAMA_ACADEMICO'].isin(equivalentes)]
    snies2 = list(programas2['CODIGO_SNIES'].unique())
    maestro2 = maestro[maestro['CODIGO_SNIES'].isin(snies2)]

    maestro3 = maestro2.merge(programas, on='CODIGO_SNIES', how='left')
    maestro4 = maestro3.merge(oferta, on=['CODIGO_SNIES', 'PERIODO'], how='left')

    return maestro4, programas2

# -------------------------------------------------------------
# Herramientas para los agentes
# -------------------------------------------------------------

@function_tool
def fetch_url(url: str, max_chars: int = 4000) -> str:
    """Descarga una página web y retorna texto visible (recortado)."""
    try:
        resp = requests.get(url, timeout=20)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        return text[:max_chars]
    except Exception as e:
        return f"Error al obtener URL: {str(e)}"

# Adaptar herramientas de mcp_tools para agents
@function_tool
def tavily_search_tool(query: str, search_depth: str = "basic") -> str:
    """Usa Tavily para búsqueda web contextual."""
    return tavily_search(query, search_depth)

@function_tool
def wikipedia_search_tool(query: str) -> str:
    """Busca información en Wikipedia."""
    result = wikipedia_search(query)
    if result.get("success"):
        return f"Título: {result.get('title', '')}\nResumen: {result.get('summary', '')}\nURL: {result.get('url', '')}"
    return f"Error: {result.get('error', 'No se encontró información')}"

@function_tool
def duckduckgo_search_tool(query: str) -> str:
    """Busca información usando DuckDuckGo."""
    result = duckduckgo_search(query)
    if result.get("success"):
        return f"Título: {result.get('title', '')}\nResumen: {result.get('summary', '')}\nURL: {result.get('url', '')}"
    return f"Error: {result.get('error', 'No se encontró información')}"

# -------------------------------------------------------------
# Agente 1: Agente Contextual
# -------------------------------------------------------------

AGENTE_CONTEXTUAL_INSTRUCTIONS = """
Eres un AGENTE CONTEXTUAL. Tu trabajo es:
1. Recibir información básica sobre un programa universitario nuevo que se quiere crear
2. Organizar esta información en un formato JSON estructurado que incluya:
   - Nombre del programa
   - Descripción
   - Nivel académico
   - Palabras clave relevantes para búsqueda
   - Contexto geográfico de interés
   - Objetivos de búsqueda específicos

El formato JSON debe ser claro y organizado para que el agente de búsqueda pueda utilizarlo eficientemente.
Debes extraer y estructurar toda la información relevante del input del usuario.
"""

# -------------------------------------------------------------
# Agente 2: Agente de Búsqueda en Línea
# -------------------------------------------------------------

AGENTE_BUSQUEDA_INSTRUCTIONS = """
Eres un AGENTE DE BÚSQUEDA EN LÍNEA. Tu trabajo es:
1. Recibir el contexto JSON del Agente Contextual
2. Buscar en la Web programas universitarios similares usando las herramientas de búsqueda
3. Extraer información de las fuentes por institución y contenido relacionado con el contexto
4. Organizar la información encontrada por programa, incluyendo:
   - Nombre del programa
   - Universidad que lo ofrece
   - País
   - URL del programa
   - Descripción
   - Cursos representativos si están disponibles
   - Fuente de la información

Busca programas a nivel local (Colombia), nacional (Latinoamérica) e internacional (EE.UU., Europa).

IMPORTANTE: Después de recopilar la información de los programas encontrados, DEBES:
1. Primero llamar a delegar_a_sumarizacion con la lista completa de programas encontrados
2. Luego, con el resultado de la sumarización, llamar a delegar_a_agregacion junto con cualquier dato SNIES disponible

Usa las herramientas de búsqueda (tavily_search, wikipedia_search, duckduckgo_search) y fetch_url cuando necesites contenido detallado.
Al final, devuelve el resultado completo del proceso de delegación.
"""

@function_tool
async def delegar_a_sumarizacion(programas_encontrados: str) -> str:
    """Delega la tarea de sumarización al agente correspondiente."""
    try:
        agente = Agent(
            name="Agente de Sumarización",
            instructions=AGENTE_SUMARIZACION_INSTRUCTIONS,
            tools=[tavily_search_tool, fetch_url, WebSearchTool()],
            model="gpt-4.1",
            model_settings=ModelSettings(temperature=0.2)
        )
        result = await Runner.run(agente, f"Resume y analiza los siguientes programas encontrados:\n\n{programas_encontrados}")
        return str(result.final_output)
    except Exception as e:
        return f"Error en sumarización: {str(e)}"

@function_tool
async def delegar_a_agregacion(resumen_y_scores: str, datos_snies: str = "") -> str:
    """Delega la tarea de agregación y generación de reporte final."""
    try:
        agente = Agent(
            name="Agente de Agregación",
            instructions=AGENTE_AGREGACION_INSTRUCTIONS,
            tools=[],
            model="gpt-4.1",
            model_settings=ModelSettings(temperature=0.3)
        )
        prompt = f"""Genera un reporte final estructurado con los siguientes datos:

RESUMEN Y SCORES:
{resumen_y_scores}

DATOS SNIES (si disponibles):
{datos_snies if datos_snies else "No hay datos SNIES disponibles para este análisis."}

El reporte debe incluir recomendaciones específicas para el nombre del programa nuevo basadas en:
- Los programas similares encontrados
- Los scores de relación
- Las tendencias en nombres identificadas
- Los datos cuantitativos disponibles"""
        result = await Runner.run(agente, prompt)
        return str(result.final_output)
    except Exception as e:
        return f"Error en agregación: {str(e)}"

# -------------------------------------------------------------
# Agente 3: Agente de Sumarización y Búsqueda
# -------------------------------------------------------------

AGENTE_SUMARIZACION_INSTRUCTIONS = """
Eres un AGENTE DE SUMARIZACIÓN Y BÚSQUEDA. Tu trabajo es:
1. Recibir la lista de programas encontrados por el Agente de Búsqueda
2. Resumir el contenido de cada programa académico
3. Analizar la descripción y el contenido de la malla curricular cuando esté disponible
4. Crear un score de relación entre cada programa encontrado y el programa original
5. Identificar tendencias en los nombres de programas similares
6. Generar un resumen comparativo

El score de relación debe ser un número entre 0 y 1, donde:
- 1.0 = programa muy similar
- 0.5-0.9 = programa relacionado
- 0.0-0.4 = programa poco relacionado

Considera factores como:
- Similitud en el nombre
- Similitud en la descripción
- Cursos comunes
- Nivel académico
- Enfoque del programa

Devuelve un resumen estructurado con los scores y las tendencias identificadas.
"""

# -------------------------------------------------------------
# Agente 4: Agente de Agregación y Reporte
# -------------------------------------------------------------

AGENTE_AGREGACION_INSTRUCTIONS = """
Eres un AGENTE DE AGREGACIÓN Y REPORTE. Tu trabajo es:
1. Recibir los resultados de la sumarización (scores y resúmenes)
2. Agrupar y organizar toda la información
3. Integrar datos cuantitativos del análisis SNIES cuando estén disponibles
4. Generar recomendaciones sobre nombres para el programa nuevo
5. Crear un reporte final estructurado que incluya:
   - Resumen ejecutivo
   - Programas similares encontrados con sus scores
   - Tendencias en nombres
   - Variables cuantitativas (número de programas, distribución geográfica, etc.)
   - Recomendaciones específicas para el nombre del programa nuevo

El reporte debe ser claro, estructurado y útil para la toma de decisiones sobre el nombre del programa.
"""

def generar_graficas(maestro4: pd.DataFrame, programa_nombre: str):
    """Genera las principales gráficas de análisis de oportunidad."""
    figuras = []

    # 1. Número de programas e instituciones en el tiempo
    NprogNies = maestro4.groupby(by='PERIODO').agg({'CODIGO_INSTITUCION_x':'nunique', 'CODIGO_SNIES':'nunique'})
    fig1, ax1 = plt.subplots()
    NprogNies.plot(ax=ax1)
    ax1.set_title(f"Programas e Instituciones en el tiempo - {programa_nombre}")
    ax1.set_xlabel('Periodo')
    ax1.set_ylabel('Cantidad')
    ax1.grid(True)
    figuras.append(fig1)

    # 2. Costo del programa vs promedio de matriculados
    maestro4['PROXY_PER'] = maestro4['PROXY_PER'].astype(int)
    df = maestro4[(maestro4['PROXY_PER']>=20211) & (maestro4['PROXY_PER']<=20242)].copy()
    df['Nombre_ies'] = df['INSTITUCION']+' - '+df['PROGRAMA_ACADEMICO']
    df = df[df['PROCESO']=='MATRICULADOS'].copy()
    df['CANTIDAD'] = df['CANTIDAD'].astype(int)
    df = df[['MATRICULA','CANTIDAD','Nombre_ies','PERIODO']].dropna()
    df = df[df['MATRICULA']!='null'].copy()
    df['MATRICULA'] = df['MATRICULA'].astype(float)
    df2 = df.groupby(by='Nombre_ies').agg({'MATRICULA':'last', 'CANTIDAD':'mean'})

    fig2, ax2 = plt.subplots()
    ax2.scatter(df2['CANTIDAD'], df2['MATRICULA'])
    for i, txt in enumerate(df2.index):
        ax2.text(df2['CANTIDAD'].iloc[i], df2['MATRICULA'].iloc[i], str(txt), fontsize=8, ha='center')
    ax2.set_xlabel('Promedio de matriculados')
    ax2.set_ylabel('Valor último de matrícula')
    ax2.set_title('Costo vs Matrícula promedio')
    ax2.grid(True)
    figuras.append(fig2)

    # 3. Valor de matrículas en el tiempo
    valor = pd.pivot_table(df, index='Nombre_ies', columns='PERIODO', values='MATRICULA', aggfunc='mean', fill_value=0)
    fig3, ax3 = plt.subplots()
    valor.T.plot(ax=ax3)
    ax3.set_title('Valor de matrícula en el tiempo')
    ax3.set_ylabel('Valor ($)')
    ax3.grid(True)
    figuras.append(fig3)

    # 4. Programas por departamento y ciudad
    df_geo = maestro4[(maestro4['PROXY_PER']>=20211) & (maestro4['PROXY_PER']<=20242)].copy()
    df_geo['Nombre_ies'] = df_geo['INSTITUCION']+' - '+df_geo['PROGRAMA_ACADEMICO']
    df_geo = df_geo[df_geo['PROCESO']=='MATRICULADOS'].copy()
    df_geo['CANTIDAD'] = df_geo['CANTIDAD'].astype(int)
    porDpto = df_geo.groupby('DEPARTAMENTO_PROGRAMA').agg({'CODIGO_SNIES':'nunique'}).sort_values(by='CODIGO_SNIES', ascending=False)
    porMpio = df_geo.groupby('MUNICIPIO_PROGRAMA').agg({'CODIGO_SNIES':'nunique'}).sort_values(by='CODIGO_SNIES', ascending=False)

    fig4, ax4 = plt.subplots(1,2, figsize=(10,4))
    porDpto.plot.bar(ax=ax4[0], legend=False, title='Programas por departamento')
    porMpio.plot.bar(ax=ax4[1], legend=False, title='Programas por municipio')
    figuras.append(fig4)

    # 5. Número de estudiantes en el tiempo
    maestro4 = maestro4[maestro4['CANTIDAD']!='null']
    maestro4['CANTIDAD'] = maestro4['CANTIDAD'].astype(int)
    num = pd.pivot_table(maestro4, index='PERIODO', columns='PROCESO', values='CANTIDAD', fill_value=0, aggfunc='sum')
    fig5, axes = plt.subplots(len(num.columns), 1, sharex=True, figsize=(8, 8))
    for i,col in enumerate(num.columns):
        axes[i].plot(num[col])
        axes[i].set_title(col)
        axes[i].grid(True)
    plt.tight_layout()
    figuras.append(fig5)

    return figuras

# -------------------------------------------------------------
# Función principal del flujo de agentes
# -------------------------------------------------------------

async def ejecutar_flujo_agentes(nombre_programa: str, descripcion: str = "", nivel: str = ""):
    """
    Ejecuta el flujo completo de agentes según el diagrama:
    1. Agente Contextual -> organiza información en JSON
    2. Agente de Búsqueda -> busca programas similares y delega
    3. Agente de Sumarización -> resume y crea scores (delegado)
    4. Agente de Agregación -> genera reporte final (delegado)
    """
    try:
        # Intentar obtener datos SNIES en paralelo (opcional)
        datos_snies = ""
        try:
            maestro4, programas2 = analizar_programa(nombre_programa)
            datos_snies = f"""
            Programas equivalentes encontrados en SNIES: {len(programas2)}
            Departamentos con presencia: {maestro4['DEPARTAMENTO_PROGRAMA'].nunique()}
            Programas únicos: {programas2['PROGRAMA_ACADEMICO'].unique().tolist()[:10]}
            """
        except:
            datos_snies = "No se pudieron obtener datos SNIES para este programa."
        
        # Paso 1: Agente Contextual
        agente_contextual = Agent(
            name="Agente Contextual",
            instructions=AGENTE_CONTEXTUAL_INSTRUCTIONS,
            tools=[],
            model="gpt-4.1",
            model_settings=ModelSettings(temperature=0.1)
        )
        
        prompt_contextual = f"""
        Programa universitario nuevo:
        - Nombre: {nombre_programa}
        - Descripción: {descripcion if descripcion else "No proporcionada"}
        - Nivel: {nivel if nivel else "No especificado"}
        
        Organiza esta información en un formato JSON estructurado con todos los campos necesarios.
        Incluye palabras clave relevantes para la búsqueda y objetivos claros.
        """
        
        resultado_contextual = await Runner.run(agente_contextual, prompt_contextual)
        contexto_json = str(resultado_contextual.final_output)
        
        # Paso 2: Agente de Búsqueda en Línea
        agente_busqueda = Agent(
            name="Agente de Búsqueda en Línea",
            instructions=AGENTE_BUSQUEDA_INSTRUCTIONS,
            tools=[
                tavily_search_tool, 
                wikipedia_search_tool, 
                duckduckgo_search_tool, 
                fetch_url, 
                delegar_a_sumarizacion, 
                delegar_a_agregacion,
                WebSearchTool()
            ],
            model="gpt-4.1",
            model_settings=ModelSettings(temperature=0.2)
        )
        
        prompt_busqueda = f"""
        Contexto del programa (JSON):
        {contexto_json}
        
        Busca programas universitarios similares en la Web usando las herramientas de búsqueda.
        Busca a nivel local (Colombia), nacional (Latinoamérica) e internacional (EE.UU., Europa).
        
        Después de encontrar y organizar la información de los programas, DEBES:
        1. Llamar a delegar_a_sumarizacion con la lista completa de programas encontrados
        2. Llamar a delegar_a_agregacion con el resultado de la sumarización y estos datos SNIES:
        {datos_snies}
        
        Devuelve el resultado final del proceso completo.
        """
        
        resultado_busqueda = await Runner.run(agente_busqueda, prompt_busqueda)
        resultado_final = str(resultado_busqueda.final_output)
        
        return {
            "contexto": contexto_json,
            "resultado_busqueda": resultado_final,
            "datos_snies": datos_snies,
            "exito": True
        }
        
    except Exception as e:
        return {
            "error": str(e),
            "exito": False
        }

# -------------------------------------------------------------
# Interfaz Streamlit
# -------------------------------------------------------------

st.title("📊 Análisis de Oportunidad de Programas Académicos (SNIES + Agentes)")
st.markdown("""
Este sistema integra el análisis de oferta académica con agentes inteligentes para buscar nombres 
de programas universitarios nuevos. El sistema utiliza un flujo de 4 agentes:
1. **Agente Contextual**: Organiza la información en formato JSON
2. **Agente de Búsqueda en Línea**: Busca programas similares en la Web
3. **Agente de Sumarización**: Resume y crea scores de relación
4. **Agente de Agregación**: Genera reporte final con recomendaciones
""")

col1, col2 = st.columns(2)

with col1:
    programa = st.text_input("Nombre del programa a analizar", value="Doctorado Ciencias Sociales")
    nivel = st.selectbox("Nivel académico", ["Pregrado", "Especialización", "Maestría", "Doctorado"], index=3)

with col2:
    descripcion = st.text_area("Descripción del programa (opcional)", 
                               value="Programa orientado a la formación de investigadores en ciencias sociales",
                               height=100)

ejecutar_agentes = st.button("🔍 Buscar nombres con Agentes", type="primary")
ejecutar_snies = st.button("📊 Análisis SNIES tradicional")

# Ejecutar flujo de agentes
if ejecutar_agentes:
    with st.spinner('🤖 Ejecutando flujo de agentes... Esto puede tomar varios minutos.'):
        try:
            resultado = asyncio.run(ejecutar_flujo_agentes(programa, descripcion, nivel))
            
            if resultado.get("exito"):
                st.success("✅ Flujo de agentes completado exitosamente")
                
                # Mostrar contexto generado
                with st.expander("📋 Contexto JSON generado por Agente Contextual", expanded=False):
                    st.code(resultado["contexto"], language="json")
                
                # Mostrar datos SNIES si están disponibles
                if resultado.get("datos_snies") and "No se pudieron obtener" not in resultado.get("datos_snies", ""):
                    with st.expander("📊 Datos SNIES integrados", expanded=False):
                        st.text(resultado["datos_snies"])
                
                # Mostrar resultado final
                st.markdown("### 📊 Reporte Final del Sistema de Agentes")
                st.markdown(resultado["resultado_busqueda"])
                
                # Intentar extraer y mostrar información estructurada
                st.markdown("---")
                st.markdown("### 💡 Recomendaciones para el nombre del programa")
                st.info("""
                El sistema de agentes ha analizado programas similares y generado recomendaciones.
                Revisa el reporte anterior para ver:
                - Programas similares encontrados
                - Scores de relación
                - Tendencias en nombres
                - Recomendaciones específicas
                """)
                
            else:
                st.error(f"❌ Error en el flujo de agentes: {resultado.get('error', 'Error desconocido')}")
                
        except Exception as e:
            st.error(f"❌ Error al ejecutar agentes: {str(e)}")
            st.exception(e)

# Ejecutar análisis SNIES tradicional
if ejecutar_snies:
    with st.spinner('Procesando información de SNIES...'):
        try:
            maestro4, programas2 = analizar_programa(programa)
            figuras = generar_graficas(maestro4, programa)

            st.success("✅ Análisis SNIES completado. Resultados:")
            st.write(f"**Programas equivalentes encontrados:** {len(programas2)}")

            for fig in figuras:
                st.pyplot(fig)

            # Crear resumen
            resumen = f"Se encontraron {len(programas2)} programas equivalentes al término '{programa}'. "\
                      f"El análisis muestra variación en matrícula y costo, con presencia en {maestro4['DEPARTAMENTO_PROGRAMA'].nunique()} departamentos."
            st.info(resumen)

            # Exportar PowerPoint
            from pptx import Presentation
            from pptx.util import Inches

            ppt_button = st.button("📄 Generar reporte PowerPoint")
            if ppt_button:
                try:
                    prs = Presentation()
                    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
                    slide_title.shapes.title.text = f"Análisis de Oportunidad - {programa}"
                    slide_title.placeholders[1].text = "Reporte generado con datos SNIES."

                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = "Resumen del Análisis"
                    slide.placeholders[1].text = resumen

                    for i, fig in enumerate(figuras, start=1):
                        buf = BytesIO()
                        fig.savefig(buf, format='png', bbox_inches='tight')
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        left, top = Inches(1), Inches(1)
                        slide.shapes.add_picture(BytesIO(buf.getvalue()), left, top, height=Inches(5)) 
                        buf.close()

                    output_path = Path('.') / f"reporte_{programa.replace(' ','_')}.pptx"
                    prs.save(str(output_path))
                    st.success(f"✅ Reporte generado: {output_path}")
                    with open(output_path, "rb") as file:
                        st.download_button(
                            label="📥 Descargar reporte PowerPoint",
                            data=file,
                            file_name=f"reporte_{programa.replace(' ','_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                except Exception as e:
                    st.warning(f"No se pudo generar el PowerPoint: {str(e)}")
                    
        except Exception as e:
            st.error(f"Error en análisis SNIES: {str(e)}")
            st.exception(e)

st.caption("Desarrollado como integración del lector SNIES con un proyecto de agentes para análisis de oportunidad académica y búsqueda de nombres de programas nuevos.")
